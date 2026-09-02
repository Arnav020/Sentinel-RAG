"""
Office loader (.docx / .pptx).

Uses python-docx and python-pptx directly rather than `unstructured`. That
dependency pulled spacy, thinc, blis, numba and llvmlite into the ingest image
to parse a handful of Office files, and it flattened everything to text with no
heading or table structure - which is exactly what the chunker needs.

Word heading styles become the heading path; slide titles do the same for
PowerPoint. Tables are emitted row-wise as atomic blocks.
"""

from __future__ import annotations

import re

from app.ingestion.document import (
    KIND_CODE,
    KIND_PROSE,
    KIND_TABLE,
    Block,
    Document,
)

_HEADING_STYLE = re.compile(r"^Heading\s*(\d)", re.IGNORECASE)
_CODEISH = re.compile(r"^(\s{2,}|\t|\$ |kubectl |apiVersion:|# )")


def _clean(s: str) -> str:
    return re.sub(r"[ \t]+", " ", s).strip()


def _rows_to_text(rows: list[list[str]]) -> str:
    return "\n".join(" | ".join(c for c in r if c) for r in rows if any(r))


def _parse_docx(file_path: str) -> Document:
    from docx import Document as DocxDocument
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = DocxDocument(file_path)
    blocks: list[Block] = []
    heading_stack: list[str] = []
    title = ""

    # Walk body children in document order so tables stay in position.
    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]

        if tag == "p":
            para = Paragraph(child, doc)
            text = _clean(para.text)
            if not text:
                continue
            style = (para.style.name or "") if para.style is not None else ""
            m = _HEADING_STYLE.match(style)
            if m or style.lower() == "title":
                level = 0 if style.lower() == "title" else int(m.group(1)) - 1
                del heading_stack[level:]
                heading_stack.append(text)
                if not title and level == 0:
                    title = text
                continue
            kind = KIND_CODE if _CODEISH.match(para.text) else KIND_PROSE
            blocks.append(Block(text=text, heading_path=tuple(heading_stack), kind=kind))

        elif tag == "tbl":
            table = Table(child, doc)
            rows = [[_clean(c.text) for c in row.cells] for row in table.rows]
            text = _rows_to_text(rows)
            if text:
                blocks.append(Block(text=text, heading_path=tuple(heading_stack), kind=KIND_TABLE))

    return Document(filename="", blocks=blocks, title=title).non_empty()


def _parse_pptx(file_path: str) -> Document:
    from pptx import Presentation

    prs = Presentation(file_path)
    blocks: list[Block] = []
    title = ""

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        if slide.shapes.title is not None:
            slide_title = _clean(slide.shapes.title.text or "")
        heading = (slide_title,) if slide_title else (f"slide {idx}",)
        if not title and slide_title:
            title = slide_title

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                rows = [[_clean(c.text) for c in row.cells] for row in shape.table.rows]
                text = _rows_to_text(rows)
                if text:
                    blocks.append(Block(text=text, heading_path=heading, kind=KIND_TABLE))
                continue
            if not shape.has_text_frame:
                continue
            parts = [_clean(p.text) for p in shape.text_frame.paragraphs if _clean(p.text)]
            if parts:
                blocks.append(Block(text="\n".join(parts), heading_path=heading, kind=KIND_PROSE))

    return Document(filename="", blocks=blocks, title=title).non_empty()


def parse_office(file_path: str) -> Document:
    lowered = file_path.lower()
    if lowered.endswith(".docx"):
        return _parse_docx(file_path)
    if lowered.endswith(".pptx"):
        return _parse_pptx(file_path)
    raise ValueError(f"Unsupported Office format: {file_path}")
